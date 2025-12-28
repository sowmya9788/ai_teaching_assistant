import os
import fitz
from typing import List, Dict, Any
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import re

class DocumentChunk:
    def __init__(self, text: str, page_num: int, doc_name: str, chunk_id: str, doc_type: str = "unknown"):
        self.text = text
        self.page_num = page_num
        self.doc_name = doc_name
        self.chunk_id = chunk_id
        self.doc_type = doc_type
        self.metadata = {
            "source": doc_name,
            "page": page_num,
            "chunk_id": chunk_id,
            "doc_type": doc_type
        }

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = ['.pdf', '.pptx', '.docx']
    
    def process_document(self, file_path: str) -> List[DocumentChunk]:
        file_ext = os.path.splitext(file_path)[1].lower()
        doc_name = os.path.basename(file_path)
        
        if file_ext == '.pdf':
            return self._process_pdf(file_path, doc_name)
        elif file_ext == '.pptx':
            return self._process_pptx(file_path, doc_name)
        elif file_ext == '.docx':
            return self._process_docx(file_path, doc_name)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
    
    def _process_pdf(self, file_path: str, doc_name: str) -> List[DocumentChunk]:
        chunks = []
        
        try:
            # Try PyMuPDF first (better for complex PDFs)
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                
                if text.strip():
                    chunk = DocumentChunk(
                        text=text.strip(),
                        page_num=page_num + 1,
                        doc_name=doc_name,
                        chunk_id=f"{doc_name}_page_{page_num + 1}",
                        doc_type="pdf"
                    )
                    chunks.append(chunk)
            doc.close()
        except:
            # Fallback to PyPDF2
            reader = PdfReader(file_path)
            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()
                if text.strip():
                    chunk = DocumentChunk(
                        text=text.strip(),
                        page_num=page_num + 1,
                        doc_name=doc_name,
                        chunk_id=f"{doc_name}_page_{page_num + 1}",
                        doc_type="pdf"
                    )
                    chunks.append(chunk)
        
        return chunks
    
    def _process_pptx(self, file_path: str, doc_name: str) -> List[DocumentChunk]:
        chunks = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                full_text = "\n".join(slide_text)
                chunk = DocumentChunk(
                    text=full_text,
                    page_num=slide_num + 1,
                    doc_name=doc_name,
                    chunk_id=f"{doc_name}_slide_{slide_num + 1}",
                    doc_type="pptx"
                )
                chunks.append(chunk)
        
        return chunks
    
    def _process_docx(self, file_path: str, doc_name: str) -> List[DocumentChunk]:
        chunks = []
        doc = Document(file_path)
        
        # Process by paragraphs
        current_page_text = []
        page_num = 1
        
        for para in doc.paragraphs:
            if para.text.strip():
                current_page_text.append(para.text.strip())
                
                # Simple page break detection
                if len(current_page_text) >= 10:  # Approximate page length
                    full_text = "\n".join(current_page_text)
                    chunk = DocumentChunk(
                        text=full_text,
                        page_num=page_num,
                        doc_name=doc_name,
                        chunk_id=f"{doc_name}_page_{page_num}",
                        doc_type="docx"
                    )
                    chunks.append(chunk)
                    current_page_text = []
                    page_num += 1
        
        # Add remaining text
        if current_page_text:
            full_text = "\n".join(current_page_text)
            chunk = DocumentChunk(
                text=full_text,
                page_num=page_num,
                doc_name=doc_name,
                chunk_id=f"{doc_name}_page_{page_num}",
                doc_type="docx"
            )
            chunks.append(chunk)
        
        return chunks
    
    def chunk_text(self, chunks: List[DocumentChunk], chunk_size: int = 500, overlap: int = 50) -> List[DocumentChunk]:
        final_chunks = []
        
        for chunk in chunks:
            text = chunk.text
            words = text.split()
            
            for i in range(0, len(words), chunk_size - overlap):
                chunk_words = words[i:i + chunk_size]
                chunk_text = " ".join(chunk_words)
                
                new_chunk = DocumentChunk(
                    text=chunk_text,
                    page_num=chunk.page_num,
                    doc_name=chunk.doc_name,
                    chunk_id=f"{chunk.chunk_id}_{i // (chunk_size - overlap)}",
                    doc_type=chunk.doc_type
                )
                final_chunks.append(new_chunk)
        
        return final_chunks
